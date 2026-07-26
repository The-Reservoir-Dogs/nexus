"""
NEXUS — Hackathon Pitch Deck generator (light, professional, 15 slides).
Run:  python3 generate_ppt.py   ->  nexus_pitch.pptx
No external assets/network needed: Pocket FM / Databricks / HackCulture are drawn as
brand-styled badges.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette
PAPER   = RGBColor(0xFA, 0xF8, 0xF4)
PANEL   = RGBColor(0xFF, 0xFD, 0xF9)
INK     = RGBColor(0x1A, 0x1A, 0x1A)
BODY    = RGBColor(0x3A, 0x37, 0x33)
MUTED   = RGBColor(0x73, 0x6D, 0x63)
LINE    = RGBColor(0xE4, 0xE0, 0xD8)
TERRA   = RGBColor(0xC2, 0x36, 0x16)   # terracotta accent
AMBER   = RGBColor(0xF3, 0xB0, 0x3A)
GREEN   = RGBColor(0x15, 0x80, 0x3D)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# brand-ish colors for badges
POCKET  = RGBColor(0xE5, 0x34, 0x4E)
DBRICKS = RGBColor(0xFF, 0x36, 0x21)
HACKC   = RGBColor(0x1A, 0x1A, 0x1A)

DISPLAY = "Georgia"      # serif, editorial
SANS    = "Calibri"      # clean body
MONO    = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


# ---------------------------------------------------------------- helpers
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    bg.shadow.inherit = False
    _send_back(bg)
    return s


def _send_back(shape):
    sp = shape._element
    sp.getparent().remove(sp)
    # insert as first child so it's behind everything
    # (python-pptx appends; move to front of spTree)
    tree = shape._element  # noqa
    return


def rect(s, x, y, w, h, color, line_color=None, line_w=0):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color
    if line_color:
        r.line.color.rgb = line_color; r.line.width = Pt(line_w)
    else:
        r.line.fill.background()
    r.shadow.inherit = False
    return r


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph = list of (txt, size, color, bold, font, italic)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, font, *rest) in para:
            italic = rest[0] if rest else False
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = font; r.font.italic = italic
    return tb


def eyebrow(s, num, label, x=Inches(0.9), y=Inches(0.62)):
    text(s, x, y, Inches(8), Inches(0.4),
         [[("// " + num + "  ", 13, TERRA, True, MONO),
           (label.upper(), 13, MUTED, True, MONO)]])


def accent_bar(s, x=Inches(0.9), y=Inches(1.15), w=Inches(0.7)):
    rect(s, x, y, w, Pt(4), TERRA)


def page_no(s, n):
    text(s, SW - Inches(1.2), SH - Inches(0.55), Inches(0.8), Inches(0.3),
         [[(f"{n:02d} / 15", 10, MUTED, False, MONO)]], align=PP_ALIGN.RIGHT)
    rect(s, Inches(0.9), SH - Inches(0.62), SW - Inches(1.8), Pt(1), LINE)


def brand_badge(s, x, y, name, color, w=Inches(1.55), h=Inches(0.42)):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = WHITE
    b.line.color.rgb = LINE; b.line.width = Pt(1)
    b.shadow.inherit = False
    tf = b.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    dot = p.add_run(); dot.text = "\u25CF "
    dot.font.size = Pt(12); dot.font.color.rgb = color; dot.font.bold = True; dot.font.name = SANS
    r = p.add_run(); r.text = name
    r.font.size = Pt(12); r.font.color.rgb = INK; r.font.bold = True; r.font.name = SANS
    return b


def partners_row(s, y):
    text(s, Inches(0.9), y - Inches(0.32), Inches(6), Inches(0.3),
         [[("POWERED BY", 10, MUTED, True, MONO)]])
    brand_badge(s, Inches(0.9), y, "Pocket FM", POCKET)
    brand_badge(s, Inches(2.65), y, "Databricks", DBRICKS)
    brand_badge(s, Inches(4.40), y, "OpenAI", INK)
    brand_badge(s, Inches(6.15), y, "HackCulture", HACKC)


def bullets(s, items, x=Inches(0.9), y=Inches(1.6), w=Inches(11.5), size=18, gap=10):
    runs = []
    for head, sub in items:
        runs.append([("\u2192  ", size, TERRA, True, SANS), (head, size, INK, True, SANS)])
        if sub:
            runs.append([("     " + sub, size - 3, BODY, False, SANS)])
    text(s, x, y, w, Inches(5), runs, space_after=gap, line_spacing=1.05)


ASSETS = os.path.join(os.path.dirname(__file__), "assets")


def img(s, name, x, y, w):
    p = os.path.join(ASSETS, name)
    if os.path.exists(p):
        pic = s.shapes.add_picture(p, x, y, width=w)
        pic.line.color.rgb = LINE; pic.line.width = Pt(1)
        return pic
    return None


def title_line(s, parts, y=Inches(1.4), x=Inches(0.9), size=40, w=Inches(11.5)):
    """parts: list of (txt, italic, color) for a single wrapped headline."""
    runs = [[(t, size, c, False, DISPLAY, it) for (t, it, c) in parts]]
    text(s, x, y, w, Inches(1.8), runs, line_spacing=1.0)


# ================================================================ SLIDES

# 1 — TITLE
s = slide()
rect(s, 0, 0, SW, Inches(0.18), TERRA)
text(s, Inches(0.9), Inches(0.9), Inches(9), Inches(0.4),
     [[("ZERO TO ONE  \u00b7  AI-NATIVE STORYTELLING (P1)", 13, MUTED, True, MONO)]])
text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.2),
     [[("NEXUS", 96, INK, True, DISPLAY)],
      [("A living story ", 30, BODY, False, DISPLAY, False),
       ("multiverse", 30, TERRA, False, DISPLAY, True),
       (".", 30, BODY, False, DISPLAY, False)]], line_spacing=1.0)
text(s, Inches(0.9), Inches(4.3), Inches(10.5), Inches(1.0),
     [[("Rewind any story, change one decision, and an AI rewrites the future \u2014", 18, BODY, False, SANS)],
      [("keeping every character and every fact consistent. Narrated as audio drama.", 18, BODY, False, SANS)]],
     line_spacing=1.15)
partners_row(s, Inches(6.2))
text(s, SW - Inches(4.3), Inches(0.9), Inches(3.4), Inches(0.4),
     [[("Team  \u00b7  The Reservoir Dogs", 13, MUTED, True, MONO)]], align=PP_ALIGN.RIGHT)

# 2 — PROBLEM
s = slide(); eyebrow(s, "01", "The problem"); accent_bar(s)
title_line(s, [("One story. ", False, INK), ("One ending.", True, TERRA)])
bullets(s, [
    ("A Pocket FM story is a single fixed path.",
     "Written once, by one team \u2014 the listener just presses play."),
    ("Personalization stops at recommendations.",
     "AI decides what you see next, but never touches the story itself."),
    ("\u201cWhat if it went differently?\u201d has no answer.",
     "The most human question a listener asks is a dead end today."),
], y=Inches(2.4))
page_no(s, 2)

# 3 — INSIGHT
s = slide(); eyebrow(s, "02", "The insight"); accent_bar(s)
title_line(s, [("Don\u2019t build an AI writer. Build a story ", False, INK),
               ("only AI makes possible.", True, TERRA)], size=36)
text(s, Inches(0.9), Inches(2.7), Inches(11.3), Inches(3),
     [[("The prompt asks what storytelling becomes possible ", 22, BODY, False, SANS),
       ("because", 22, TERRA, True, SANS),
       (" AI exists.", 22, BODY, False, SANS)],
      [("", 10, BODY, False, SANS)],
      [("Answer: a universe of infinite branching timelines that a single AI keeps ", 22, INK, False, SANS),
       ("perfectly consistent", 22, TERRA, True, SANS),
       (" \u2014 something no human writers\u2019 room could ever maintain by hand.", 22, INK, False, SANS)]],
     line_spacing=1.2, space_after=8)
page_no(s, 3)

# 4 — SOLUTION
s = slide(); eyebrow(s, "03", "The solution"); accent_bar(s)
title_line(s, [("NEXUS \u2014 the story ", False, INK), ("time machine.", True, TERRA)])
bullets(s, [
    ("Rewind to any decision point in an episode.",
     "\u201cThe hero spared the villain\u201d \u2192 what if she didn\u2019t?"),
    ("AI regenerates the alternate future \u2014 consistently.",
     "Remembers each character\u2019s voice + what they can know."),
    ("Plays back as a cinematic multi-voice audio drama.",
     "Pocket FM\u2019s native format \u2014 not text on a screen."),
    ("Community rates; the author verifies canon.",
     "One story becomes infinite, coherent timelines."),
], y=Inches(2.4), w=Inches(6.2), gap=9)
img(s, "multiverse.png", Inches(7.35), Inches(2.35), Inches(5.4))
page_no(s, 4)

# 5 — HOW IT WORKS
s = slide(); eyebrow(s, "04", "How it works"); accent_bar(s)
title_line(s, [("From one decision to a new ", False, INK), ("consistent world.", True, TERRA)], size=34)
steps = ["Rewind\na decision", "AI gathers\ncanon + memory", "Regenerate\nthe timeline",
         "Self-check\nconsistency", "Narrate as\naudio drama", "Community\ncanonizes"]
n = len(steps); bw = Inches(1.75); gap = Inches(0.22)
total = Emu(int(bw) * n + int(gap) * (n - 1))
x0 = Emu(int((SW - total) / 2)); y0 = Inches(3.0)
for i, st in enumerate(steps):
    x = Emu(int(x0) + i * (int(bw) + int(gap)))
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y0, bw, Inches(1.7))
    card.fill.solid(); card.fill.fore_color.rgb = PANEL
    card.line.color.rgb = LINE; card.line.width = Pt(1); card.shadow.inherit = False
    tf = card.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    rr = p.add_run(); rr.text = f"{i+1}\n"; rr.font.size = Pt(20); rr.font.bold = True
    rr.font.color.rgb = TERRA; rr.font.name = DISPLAY
    for j, ln in enumerate(st.split("\n")):
        pp = tf.add_paragraph(); pp.alignment = PP_ALIGN.CENTER
        r = pp.add_run(); r.text = ln; r.font.size = Pt(12); r.font.color.rgb = INK; r.font.name = SANS
    if i < n - 1:
        text(s, Emu(int(x) + int(bw) - Inches(0.05)), Emu(int(y0) + int(Inches(0.65))),
             Inches(0.3), Inches(0.4), [[("\u2192", 16, TERRA, True, SANS)]], align=PP_ALIGN.CENTER)
page_no(s, 5)

# 6 — FEATURE: Time machine + memory
s = slide(); eyebrow(s, "05", "Key feature \u00b7 01"); accent_bar(s)
title_line(s, [("Story Time Machine + ", False, INK), ("persistent memory.", True, TERRA)], size=34)
bullets(s, [
    ("Change one decision, regenerate every future.",
     "The core P1 challenge \u2014 solved, not simulated."),
    ("Characters keep voice, goals + knowledge limits.",
     "Never know what didn\u2019t happen in this timeline."),
    ("Each side character carries its own memory.",
     "The universe grows without breaking continuity."),
], y=Inches(2.55), w=Inches(6.0))
img(s, "splitview.png", Inches(7.15), Inches(2.5), Inches(5.6))
page_no(s, 6)

# 7 — FEATURE: Self-proving consistency (WOW)
s = slide()
rect(s, 0, 0, SW, SH, PAPER)
rect(s, 0, 0, Inches(0.22), SH, TERRA)
eyebrow(s, "06", "Key feature \u00b7 02  \u2014  our wow", x=Inches(1.1))
title_line(s, [("The AI ", False, INK), ("grades and repairs itself", True, TERRA),
               (" \u2014 live.", False, INK)], size=34, x=Inches(1.1))
bullets(s, [
    ("Scored: continuity, character, reader-intent.",
     "With exact quoted evidence \u2014 not \u2018looks good\u2019."),
    ("Catches a contradiction? It fixes itself, live.",
     "Watch the score climb in front of the jury."),
    ("Nothing unsafe auto-publishes \u2014 a human decides.",
     "Eval + human-in-the-loop, one story."),
], x=Inches(1.1), y=Inches(2.55), w=Inches(5.3))
img(s, "eval.png", Inches(6.85), Inches(2.15), Inches(5.9))
text(s, Inches(1.1), Inches(6.55), Inches(11), Inches(0.6),
     [[("Every team claims \u201cconsistent.\u201d We ", 16, INK, False, SANS),
       ("prove it \u2014 and self-correct when wrong.", 16, TERRA, True, SANS)]])
page_no(s, 7)

# 8 — FEATURE: Cinematic audio
s = slide(); eyebrow(s, "07", "Key feature \u00b7 03"); accent_bar(s)
title_line(s, [("Cinematic ", False, INK), ("audio drama", True, TERRA),
               (", not text.", False, INK)], size=34)
bullets(s, [
    ("Distinct multi-voice narration per branch.", ""),
    ("Music bed + a beat of silence before the moment.", ""),
    ("Audio-first \u2014 a layer on Pocket FM\u2019s strength.", ""),
], y=Inches(2.35), gap=6)
img(s, "audio.png", Inches(1.6), Inches(3.75), Inches(10.1))
page_no(s, 8)

# 9 — FEATURE: Co-author + community
s = slide(); eyebrow(s, "08", "Key feature \u00b7 04"); accent_bar(s)
title_line(s, [("AI Co-Author + ", False, INK), ("community canon.", True, TERRA)], size=34)
bullets(s, [
    ("Readers steer the story with comments and ratings.",
     "Their signals become the input that drives the next branch."),
    ("Co-authors write with an AI editor \u2014 human approves every line.",
     "A VS-Code-style studio; nothing saves until the human signs off."),
    ("The original author verifies which branch becomes canon.",
     "The universe grows, curated \u2014 AI keeps it coherent."),
], y=Inches(2.5))
page_no(s, 9)

# 10 — LIVE DEMO
s = slide()
rect(s, 0, 0, SW, SH, INK)
text(s, Inches(0.9), Inches(0.7), Inches(9), Inches(0.4),
     [[("// LIVE DEMO", 14, AMBER, True, MONO)]])
text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(2),
     [[("Let\u2019s rewrite a fate.", 60, WHITE, False, DISPLAY, True)]])
text(s, Inches(0.9), Inches(4.4), Inches(11), Inches(1),
     [[("Rewind \u2192 regenerate \u2192 watch it catch & fix itself \u2192 hear the branch.", 20,
        RGBColor(0xD6, 0xCF, 0xC1), False, SANS)]])
page_no_dark = text(s, SW - Inches(1.4), SH - Inches(0.55), Inches(1), Inches(0.3),
                    [[("10 / 15", 10, RGBColor(0x9A, 0x9A, 0x9A), False, MONO)]], align=PP_ALIGN.RIGHT)

# 11 — ARCHITECTURE
s = slide(); eyebrow(s, "09", "Under the hood"); accent_bar(s)
title_line(s, [("Built ", False, INK), ("Databricks-native.", True, TERRA)], size=34)
comp = [("Databricks Apps", "Serverless hosting + OAuth for the web app & agent"),
        ("Lakebase (Postgres)", "Single store: series, episodes, timelines, character memory"),
        ("Foundation Model APIs", "The tool-calling AI co-author that regenerates timelines"),
        ("Mosaic AI Agent", "Reads canon via tools, streams its reasoning live"),
        ("MLflow 3 GenAI Eval", "The consistency judges + the self-correction scores"),
        ("Multi-voice TTS", "Renders each branch into an audio drama")]
colw = Inches(5.6); rowh = Inches(0.95)
for i, (h, d) in enumerate(comp):
    col = i % 2; row = i // 2
    x = Inches(0.9) + col * (colw + Inches(0.3)); y = Inches(2.5) + row * (rowh + Inches(0.12))
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, colw, rowh)
    card.fill.solid(); card.fill.fore_color.rgb = PANEL
    card.line.color.rgb = LINE; card.line.width = Pt(1); card.shadow.inherit = False
    text(s, x + Inches(0.25), y + Inches(0.12), colw - Inches(0.5), rowh,
         [[(h, 15, TERRA, True, SANS)], [(d, 12, BODY, False, SANS)]], space_after=2, line_spacing=1.0)
page_no(s, 11)

# 12 — DIFFERENTIATION
s = slide(); eyebrow(s, "10", "Why we\u2019re different"); accent_bar(s)
title_line(s, [("Not \u201cAI writes a story.\u201d", True, INK)], size=34)
# two columns compare
lx, rx = Inches(0.9), Inches(6.95); cy = Inches(2.5); cw = Inches(5.6)
def compare_col(x, label, color, items):
    hdr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, cy, cw, Inches(0.55))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = color; hdr.line.fill.background(); hdr.shadow.inherit = False
    tf = hdr.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.size = Pt(15); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = SANS
    runs = [[("\u2022  " + it, 16, INK if color == TERRA else BODY, False, SANS)] for it in items]
    text(s, x + Inches(0.15), cy + Inches(0.75), cw - Inches(0.3), Inches(3.4), runs,
         space_after=10, line_spacing=1.05)
compare_col(lx, "OTHER TEAMS", MUTED,
            ["One prompt \u2192 one story", "\u2018Trust us it\u2019s consistent\u2019",
             "Generic single-voice TTS", "A one-shot demo"])
compare_col(rx, "NEXUS", TERRA,
            ["A universe of consistent timelines", "Proves consistency \u2014 and self-corrects",
             "Multi-voice cinematic audio", "A platform Pocket FM could ship"])
page_no(s, 12)

# 13 — IMPACT
s = slide(); eyebrow(s, "11", "Why it matters for Pocket FM"); accent_bar(s)
title_line(s, [("Infinite content. Deeper ", False, INK), ("engagement.", True, TERRA)], size=34)
stats = [("\u221e", "Timelines from one story", "More catalogue, zero extra writers\u2019 rooms"),
         ("+", "Listener retention", "\u201cMy story\u201d beats \u201ca story\u201d \u2014 people finish what feels theirs"),
         ("\u25CE", "A new format", "Interactive, personal audio drama \u2014 AI-native, on-brand")]
cw = Inches(3.7)
for i, (big, h, d) in enumerate(stats):
    x = Inches(0.9) + i * (cw + Inches(0.3)); y = Inches(2.7)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, Inches(3.0))
    card.fill.solid(); card.fill.fore_color.rgb = PANEL
    card.line.color.rgb = LINE; card.line.width = Pt(1); card.shadow.inherit = False
    text(s, x + Inches(0.3), y + Inches(0.3), cw - Inches(0.6), Inches(1.2),
         [[(big, 54, TERRA, True, DISPLAY)]])
    text(s, x + Inches(0.3), y + Inches(1.5), cw - Inches(0.6), Inches(1.4),
         [[(h, 18, INK, True, SANS)], [(d, 13, BODY, False, SANS)]], space_after=6, line_spacing=1.05)
page_no(s, 13)

# 14 — ROADMAP
s = slide(); eyebrow(s, "12", "What\u2019s next"); accent_bar(s)
title_line(s, [("From hackathon to ", False, INK), ("product.", True, TERRA)], size=34)
bullets(s, [
    ("Now  \u2014  live demo", "Rewind \u2192 consistent regeneration \u2192 self-correcting evals \u2192 audio drama."),
    ("Next  \u2014  scale the universe", "Vector memory for deep lineage; multi-level branching."),
    ("Then  \u2014  Pocket FM integration", "A personalization layer over the existing catalogue & voice pipeline."),
    ("Vision  \u2014  every side character a protagonist", "Infinite Story Universe: today\u2019s villain is tomorrow\u2019s hero."),
], y=Inches(2.5), gap=9)
page_no(s, 14)

# 15 — CLOSING
s = slide()
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, 0, SW, Inches(0.18), TERRA)
text(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(2.2),
     [[("NEXUS", 84, WHITE, True, DISPLAY)],
      [("Every story has infinite endings. ", 26, RGBColor(0xD6, 0xCF, 0xC1), False, DISPLAY, False),
       ("We keep them true.", 26, AMBER, False, DISPLAY, True)]], line_spacing=1.0)
text(s, Inches(0.9), Inches(4.5), Inches(11), Inches(1),
     [[("Other teams\u2019 AI claims to be consistent.", 20, WHITE, False, SANS)],
      [("Ours ", 20, WHITE, False, SANS), ("proves it \u2014 and fixes itself when it\u2019s wrong.", 20, AMBER, True, SANS)]],
     line_spacing=1.2)
# partner badges on dark
for i, (nm, col) in enumerate([("Pocket FM", POCKET), ("Databricks", DBRICKS), ("OpenAI", WHITE), ("HackCulture", WHITE)]):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9) + i * Inches(1.75), Inches(6.2), Inches(1.55), Inches(0.42))
    b.fill.background(); b.line.color.rgb = RGBColor(0x44, 0x40, 0x3A); b.line.width = Pt(1); b.shadow.inherit = False
    tf = b.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    d = p.add_run(); d.text = "\u25CF "; d.font.size = Pt(11); d.font.color.rgb = col; d.font.bold = True; d.font.name = SANS
    r = p.add_run(); r.text = nm; r.font.size = Pt(11); r.font.color.rgb = WHITE; r.font.name = SANS
text(s, SW - Inches(4.5), Inches(6.25), Inches(3.6), Inches(0.4),
     [[("Team \u00b7 The Reservoir Dogs", 12, RGBColor(0x9A, 0x9A, 0x9A), True, MONO)]], align=PP_ALIGN.RIGHT)

prs.save("nexus_pitch.pptx")
print("saved nexus_pitch.pptx  \u2014  slides:", len(prs.slides._sldIdLst))
